import html
import pathlib
import re
from io import BytesIO
from typing import Dict, Union, Optional
try:
    import pptx
except:
    raise ImportError("`python-pptx` not detected, please install"
                      "using `pip install msgflow[python-pptx]`")
from msgflow.data.parsers.base import BaseParser
from msgflow.data.parsers.types import PptxParser


class PythonPPTXPptxParser(BaseParser, PptxParser):

    """ Python-PPTX Pptx Parser is a module to convert 
    .pptx in markdown. 
    
    Python-PPTX markdown parser is based-on Markitdown.

    This module is able to extract images and return them
    as BytesIO (BufferReader).
    """
    
    provider = "python_pptx"

    def __init__(self):
        pass

    def __call__(self, path: str) -> Dict[str, Union[str, Dict[str, BytesIO]]]:
        if pathlib.Path(path).suffix.lower() == ".pptx":
            return self._convert(path)
        else:
            ValueError("`Python-PPTX` requires a path that "
                       f"ends with `.pptx`, given `{path}`")

    def _convert(self, path):
        md_content = ""
        images_dict = {}

        presentation = pptx.Presentation(path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # Extract image data
                    image_buffer = self._extract_image_data(shape)
                    if image_buffer:
                        # Generate unique filename
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        
                        # Store image data in dictionary
                        images_dict[filename] = image_buffer

                        # Get alt text
                        alt_text = ""
                        try:
                            alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                        except Exception:
                            pass

                        md_content += (
                            "\n!["
                            + (alt_text if alt_text else shape.name)
                            + "]("
                            + filename
                            + ")\n"
                        )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return {
            "text": md_content.strip(),
            "images": images_dict
        }

    def _extract_image_data(self, shape) -> Optional[bytes]:
        """ Extract image data from a shape as BytesIO object """
        try:
            if hasattr(shape, "image"):
                # For placeholder shapes with images
                return BytesIO(shape.image.blob)
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                # For regular picture shapes
                return BytesIO(shape._element.blip.blob)
        except Exception:
            return None
        return None

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_chart_to_markdown(self, chart):
        md = "\n\n### Chart"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"
        data = []
        category_names = [c.label for c in chart.plots[0].categories]
        series_names = [s.name for s in chart.series]
        data.append(["Category"] + series_names)

        for idx, category in enumerate(category_names):
            row = [category]
            for series in chart.series:
                row.append(series.values[idx])
            data.append(row)

        markdown_table = []
        for row in data:
            markdown_table.append("| " + " | ".join(map(str, row)) + " |")
        header = markdown_table[0]
        separator = "|" + "|".join(["---"] * len(data[0])) + "|"
        return md + "\n".join([header, separator] + markdown_table[1:])